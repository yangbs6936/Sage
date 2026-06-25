"""
PPTX文件解析器
支持PowerPoint文件的文本提取和元数据获取
"""

import traceback
import os
import subprocess
from typing import Dict, Any
from pptx import Presentation
from .base_parser import BaseFileParser, ParseResult
from sagents.utils.logger import logger


class PPTXParser(BaseFileParser):
    """PPTX文件解析器"""

    SUPPORTED_EXTENSIONS = [".pptx", ".ppt"]
    SUPPORTED_MIME_TYPES = [
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    ]

    def parse(self, file_path: str, skip_validation: bool = False) -> ParseResult:
        """
        解析PPTX文件

        Args:
            file_path: PPTX文件路径
            skip_validation: 是否跳过文件格式验证（can_parse检查）

        Returns:
            ParseResult: 解析结果
        """
        if not self.validate_file(file_path):
            return self.create_error_result(
                f"文件不存在或无法读取: {file_path}", file_path
            )

        if not skip_validation and not self.can_parse(file_path):
            return self.create_error_result(f"不支持的文件类型: {file_path}", file_path)

        temp_pptx_path = None

        try:
            # 检查文件扩展名
            ext = os.path.splitext(file_path)[1].lower()

            # 如果是PPT文件，先转换为PPTX
            if ext == ".ppt":
                try:
                    temp_pptx_path = self._convert_ppt_to_pptx(file_path)
                    target_path = temp_pptx_path
                except Exception as e:
                    return self.create_error_result(f"PPT转换失败: {str(e)}", file_path)
            else:
                target_path = file_path

            # 加载演示文稿
            prs = Presentation(target_path)

            # 提取文本内容
            text_parts = []
            slide_texts = []

            for i, slide in enumerate(prs.slides, 1):
                slide_text = f"--- 幻灯片 {i} ---\\n"
                slide_content = []

                # 提取幻灯片中的文本
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():  # pyright: ignore[reportAttributeAccessIssue]
                        slide_content.append(shape.text.strip())  # pyright: ignore[reportAttributeAccessIssue]
                    # 处理表格
                    if shape.has_table:
                        table = shape.table  # pyright: ignore[reportAttributeAccessIssue]
                        table_text = []
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                table_text.append(" | ".join(row_text))
                        if table_text:
                            slide_content.append("\\n".join(table_text))

                if slide_content:
                    slide_text += "\\n".join(slide_content)
                else:
                    slide_text += "(无文本内容)"

                slide_texts.append("\\n".join(slide_content))
                text_parts.append(slide_text)

            text = "\\n\\n".join(text_parts)

            # 获取基础文件元数据
            base_metadata = self.get_file_metadata(file_path)

            # 获取PPTX特定元数据
            pptx_metadata = self._extract_pptx_metadata(prs, slide_texts)

            # 合并元数据
            metadata = {**base_metadata, **pptx_metadata}

            # 添加文本统计信息
            metadata.update(
                {
                    "text_length": len(text),
                    "character_count": len(text),
                    "word_count": len(text.split()) if text else 0,
                    "line_count": text.count("\\n") if text else 0,
                }
            )

            return ParseResult(text=text, metadata=metadata, success=True)

        except Exception as e:
            error_msg = f"PPTX解析失败: {str(e)}"
            logger.error(error_msg)
            traceback.print_exc()
            return self.create_error_result(error_msg, file_path)
        finally:
            # 清理临时文件
            if temp_pptx_path and os.path.exists(temp_pptx_path):
                try:
                    os.remove(temp_pptx_path)
                except Exception:
                    pass

    def _convert_ppt_to_pptx(self, file_path: str) -> str:
        """
        使用LibreOffice将PPT转换为PPTX
        """
        pptx_output_path = file_path + "x"

        # 检查libreoffice是否安装
        try:
            subprocess.run(
                ["libreoffice", "--version"], check=True, capture_output=True
            )
        except (FileNotFoundError, subprocess.CalledProcessError):
            raise Exception("LibreOffice未安装或不在PATH中，无法转换PPT文件")

        logger.info(
            f"尝试使用LibreOffice将PPT转换为PPTX: {file_path} -> {pptx_output_path}"
        )

        command = [
            "libreoffice",
            "--headless",
            "--convert-to",
            "pptx",
            "--outdir",
            os.path.dirname(file_path),
            file_path,
        ]

        result = subprocess.run(command, capture_output=True, text=True)

        if result.returncode != 0:
            logger.error(f"LibreOffice转换错误: {result.stderr}")
            raise Exception(f"LibreOffice转换失败: {result.stderr}")

        if not os.path.exists(pptx_output_path):
            # 有时候输出文件名可能不同，尝试根据文件名猜测
            # 这里简单处理，假设如果上面的check通过，文件应该存在
            raise Exception(f"转换后的文件未找到: {pptx_output_path}")

        return pptx_output_path

    def _extract_pptx_metadata(self, prs: Any, slide_texts: list) -> Dict[str, Any]:
        """
        提取PPTX特定元数据

        Args:
            prs: Presentation对象
            slide_texts: 幻灯片文本列表

        Returns:
            Dict[str, Any]: PPTX元数据
        """
        try:
            # 获取文档属性
            core_props = prs.core_properties

            # 计算统计信息
            slide_count = len(prs.slides)
            total_shapes = sum(len(slide.shapes) for slide in prs.slides)
            text_shapes = 0
            image_shapes = 0

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_shapes += 1
                    if hasattr(shape, "image"):
                        image_shapes += 1

            # 计算每张幻灯片的文本长度
            slide_text_lengths = [len(text) for text in slide_texts]

            metadata = {
                # 文档属性
                "title": getattr(core_props, "title", "") or "",
                "author": getattr(core_props, "author", "") or "",
                "subject": getattr(core_props, "subject", "") or "",
                "keywords": getattr(core_props, "keywords", "") or "",
                "comments": getattr(core_props, "comments", "") or "",
                "category": getattr(core_props, "category", "") or "",
                "created": str(getattr(core_props, "created", ""))
                if getattr(core_props, "created", None)
                else "",
                "modified": str(getattr(core_props, "modified", ""))
                if getattr(core_props, "modified", None)
                else "",
                "last_modified_by": getattr(core_props, "last_modified_by", "") or "",
                "revision": getattr(core_props, "revision", 0) or 0,
                # 演示文稿统计
                "slide_count": slide_count,
                "total_shapes": total_shapes,
                "text_shapes": text_shapes,
                "image_shapes": image_shapes,
                "slide_text_lengths": slide_text_lengths,
                "average_text_per_slide": sum(slide_text_lengths) / slide_count
                if slide_count > 0
                else 0,
                # 幻灯片尺寸
                "slide_width": prs.slide_width,
                "slide_height": prs.slide_height,
                "slide_size_inches": {
                    "width": prs.slide_width / 914400,  # 转换为英寸
                    "height": prs.slide_height / 914400,
                },
            }

            return metadata

        except Exception as e:
            print(f"提取PPTX元数据时出错: {e}")
            traceback.print_exc()
            return {"metadata_extraction_error": str(e)}
